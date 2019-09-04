from pptx import Presentation
import datetime
import tempfile

class ClinicalPatientPPT():
    def __init__(self, patientid, clinical_center, username):
        self.patientid = patientid
        self.clinical_center = clinical_center
        self.username = username

        # additional metadata upon creation
        self.date = datetime.datetime.now()

        # initialize presentation
        self.prs = Presentation()
        self.tempfile = tempfile.TemporaryFile(mode="wb")

    @property
    def cached_file_(self):
        return self.tempfile

    def _save(self):
        self.prs.save(self.tempfile)

    def create_title_slide(self, titletxt="Hello, World!", subtitletxt="python-pptx was here!"):
        title_slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = titletxt
        subtitle.text = subtitletxt
        self._save()

    def add_slide(self, title, content, photos):
        pass

    def add_clinical_summary(self, patient_data_list: list):
        for i, pdata in enumerate(patient_data_list):
            # add a slide w/ corresponding data
            self.add_slide(title, content, photos)
            bullet_slide_layout = prs.slide_layouts[1]

            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes

            title_shape = shapes.title
            body_shape = shapes.placeholders[1]

            title_shape.text = 'Adding a Bullet Slide'

            tf = body_shape.text_frame
            tf.text = 'Find the bullet slide layout'

            p = tf.add_paragraph()
            p.text = 'Use _TextFrame.text for first bullet'
            p.level = 1

            p = tf.add_paragraph()
            p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
            p.level = 2

